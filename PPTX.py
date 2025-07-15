from fastapi import FastAPI, File, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from pptx import Presentation
import tempfile
import json
import re
import uvicorn
import os
import subprocess
import base64

# Google Generative AI
import google.generativeai as genai

app = FastAPI()

# Allow CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

# Configure Gemini
genai.configure(api_key="get your API key from https://console.cloud.google.com/apis/credentials")
model = genai.GenerativeModel("gemini-1.5-flash")


def normalize_json(text: str) -> str:
    """Clean invalid JSON characters."""
    text = text.replace("\u201c", '"').replace("\u201d", '"') \
               .replace("\u2018", "'").replace("\u2019", "'")
    text = re.sub(r"(?<!\\)\\(?![\"/bfnrtu])", "", text)
    return text


def extract_json_from_text(text: str):
    """Extract JSON block from model response."""
    match = re.search(r'\[.*\]', text, re.DOTALL)  # Expecting a JSON array
    if match:
        cleaned = normalize_json(match.group(0))
        return json.loads(cleaned)
    else:
        raise ValueError("No JSON array found in Gemini response")


def chunk_list(lst, chunk_size):
    """Yield successive chunk_size-sized chunks from lst."""
    for i in range(0, len(lst), chunk_size):
        yield lst[i : i + chunk_size]


def generate_quiz_batch(slides_chunk, settings=None):
    """
    Generate quiz content for a batch of slides.
    Expects slides_chunk to be a list of dictionaries each with "slide", "text", and "speaker_notes".
    Optionally uses settings (totalQuestions, hardCount, mediumCount, easyCount, allowHints, allowAIAsk, quizNotes)
    to guide the quiz generation.
    Returns a list of quiz objects including a hint.
    """
    # Log each slide's data for debugging
    for slide in slides_chunk:
        print(f"Processing slide {slide.get('slide')}: {slide.get('text')[:60]}...")

    # Build the prompt
    prompt = (
        "You are an AI quiz generator meant to receive PPTX slides and create explanations and questions to guide the student. "
        "Do not refer to each slide individually but treat this as content for a course. "
        "Create multiple-choice quiz questions for the following slides. For each slide, use the provided text and speaker notes to generate a quiz question with four answer choices (A, B, C, D), the correct answer (a letter), an explanation, and a subtle hint for the student. "
        "Format your response as a JSON array of objects. Each object should have the keys:\n"
        '"slide": <slide number>, "question": <string>, "choices": [<choiceA>, <choiceB>, <choiceC>, <choiceD>], "answer": <letter>, "explanation": <string>, "hint": <string>.\n'
        "If a slide has nothing to analyze then you may return the same structure with 'N/A' in everything.\n\n"
    )
    # Add quiz settings if present
    if settings:
        prompt += (
            "Quiz Settings:\n"
            f"Total Questions: {settings.get('totalQuestions', 'N/A')}\n"
            f"Hard Questions: {settings.get('hardCount', 'N/A')}\n"
            f"Medium Questions: {settings.get('mediumCount', 'N/A')}\n"
            f"Easy Questions: {settings.get('easyCount', 'N/A')}\n"
            f"Allow Hints: {settings.get('allowHints', 'N/A')}\n"
            f"Allow AI Ask: {settings.get('allowAIAsk', 'N/A')}\n"
            f"Quiz Notes: {settings.get('quizNotes', 'N/A')}\n\n"
        )

    prompt += "Here are the slides:\n"
    for slide in slides_chunk:
        slide_num = slide.get("slide")
        slide_text = slide.get("text", "")
        speaker_notes = slide.get("speaker_notes", "N/A")
        prompt += f"Slide {slide_num}: {slide_text}\nSpeaker Notes: {speaker_notes}\n\n"

    try:
        response = model.generate_content(prompt)
        text_response = response.text.strip()
        quiz_data = extract_json_from_text(text_response)  # Expecting a JSON array
        return quiz_data
    except Exception as e:
        print(f"Error generating quiz for batch starting at slide {slides_chunk[0].get('slide')}: {e}")
        # Return default placeholders if an error occurs
        default_quizzes = []
        for slide in slides_chunk:
            slide_num = slide.get("slide")
            default_quizzes.append({
                "slide": slide_num,
                "question": "No valid question found.",
                "choices": ["A) ...", "B) ...", "C) ...", "D) ..."],
                "answer": "A",
                "explanation": f"No explanation available for slide {slide_num}.",
                "hint": "No hint available."
            })
        return default_quizzes


# --- Slide Extraction using PowerPoint COM ---
import win32com.client

def convert_pptx_to_images(pptx_path: str, output_dir: str):
    """Converts a PPTX to slide images (PNG) using PowerPoint via COM Automation."""
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    powerpoint.Visible = 1

    try:
        # Open PPTX file
        presentation = powerpoint.Presentations.Open(pptx_path, WithWindow=False)

        # Create output folder if not exists
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

        # Export each slide to PNG
        for i, slide in enumerate(presentation.Slides):
            image_path = os.path.join(output_dir, f"Slide{i + 1}.png")
            slide.Export(image_path, "PNG")

        presentation.Close()
    except Exception as e:
        print(f"Error during PPTX to image conversion: {e}")
    finally:
        powerpoint.Quit()


@app.post("/extract_slides/")
async def extract_slides(file: UploadFile = File(...)):
    """Extract text, speaker notes, and base64 slide images from PPTX."""
    with tempfile.TemporaryDirectory() as tmpdir:
        tmp_pptx = os.path.join(tmpdir, "upload.pptx")
        with open(tmp_pptx, "wb") as f:
            f.write(await file.read())

        # Convert slides to images
        convert_pptx_to_images(tmp_pptx, tmpdir)

        # Parse text and notes
        prs = Presentation(tmp_pptx)
        slides = []
        for i, slide in enumerate(prs.slides):
            text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
            speaker_notes = ""
            if hasattr(slide, "has_notes_slide") and slide.has_notes_slide:
                try:
                    speaker_notes = slide.notes_slide.notes_text_frame.text.strip()
                except Exception as e:
                    print(f"Error retrieving speaker notes for slide {i+1}: {e}")
                    speaker_notes = "N/A"
            image_path = os.path.join(tmpdir, f"Slide{i+1}.png")

            image_b64 = None
            if os.path.exists(image_path):
                with open(image_path, "rb") as img_file:
                    image_b64 = base64.b64encode(img_file.read()).decode("utf-8")

            slides.append({
                "slide": i + 1,
                "text": text.strip(),
                "speaker_notes": speaker_notes,
                "image_base64": image_b64
            })

        print(f"Extracted {len(slides)} slide(s)")
        return {"slides": slides}


@app.post("/quiz/")
async def generate_quiz_endpoint(payload: dict):
    """Generate quiz content for selected slides in batches of 15."""
    print("Raw payload received for quiz generation:")
    print(json.dumps(payload, indent=2))

    selected = payload.get("selected", [])
    if not selected:
        print("No slides selected for quiz generation.")
        return {"sections": []}

    # Remove 'image_base64' from each slide to avoid large payload
    for slide in selected:
        if "image_base64" in slide:
            del slide["image_base64"]

    settings = payload.get("settings", {})
    for slide in selected:
        print(f"Selected slide {slide.get('slide')}: {slide.get('text')[:60]}...")

    results = []
    # Process in batches of 15
    for chunk in chunk_list(selected, 15):
        batch_result = generate_quiz_batch(chunk, settings)
        results.extend(batch_result)

    final_output = {"sections": results}
    print("Sending quiz results:")
    print(json.dumps(final_output, indent=2))
    return final_output


@app.post("/ask/")
async def ask_question(payload: dict):
    """
    Given the original slide details along with a quiz question and its content,
    answer the student's follow-up question.

    Expected payload keys:
      - slide: slide number
      - text: original slide text
      - question: the quiz question text
      - answer: the correct answer
      - explanation: explanation for the quiz question
      - student_question: the student's follow-up query

    Returns a concise, step-by-step guided answer.
    """
    prompt = (
        "You are an AI assistant. Based on the slide details and the quiz question provided below, "
        "answer the student's follow-up question in a clear and comprehensive manner.\n\n"
        "Slide: {slide}\n"
        "Slide Text: {text}\n"
        "Quiz Question: {question}\n"
        "Correct Answer: {answer}\n"
        "Explanation: {explanation}\n"
        "Student Question: {student_question}\n\n"
        "Please provide a simple, short, step-by-step answer that ties together the slide content and quiz context."
    ).format(
        slide=payload.get("slide", "N/A"),
        text=payload.get("text", "N/A"),
        question=payload.get("question", "N/A"),
        answer=payload.get("answer", "N/A"),
        explanation=payload.get("explanation", "N/A"),
        student_question=payload.get("student_question", "N/A")
    )
    try:
        response = model.generate_content(prompt)
        return {"response": response.text.strip()}
    except Exception as e:
        return {"error": str(e)}


if __name__ == "__main__":
    uvicorn.run(app, host="127.0.0.1", port=8000)
